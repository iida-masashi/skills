"""Smoke tests for deliverable-review skill.

これらのテストは外部API（Gemini）を叩きません。
LLM レビューのテストは test_llm_review.py で別途行う。
"""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path

import pytest

SKILL_ROOT = Path(__file__).resolve().parent.parent
SCRIPTS_DIR = SKILL_ROOT / "scripts"

sys.path.insert(0, str(SCRIPTS_DIR))


@pytest.fixture(scope="module")
def sample_pptx(tmp_path_factory):
    """_smoke_make_sample.py で生成した pptx を返す。"""
    tmp = tmp_path_factory.mktemp("sample")
    out = tmp / "sample.pptx"
    subprocess.run(
        [sys.executable, str(SCRIPTS_DIR / "_smoke_make_sample.py"), str(out)],
        check=True,
    )
    assert out.exists()
    return out


def test_renamed_module_importable():
    """リネーム後のモジュールが import できる。"""
    import ai_check_extract
    import llm_review
    import strategy_checks

    assert callable(ai_check_extract.write_ai_check_json)
    assert callable(ai_check_extract.write_prompt_hint)
    assert callable(llm_review.run_llm_review)
    assert callable(strategy_checks.run_strategy_checks)


def test_old_module_absent():
    """旧 phase3_extract.py が削除されている。"""
    assert not (SCRIPTS_DIR / "phase3_extract.py").exists()


def test_ai_check_json_output(sample_pptx, tmp_path):
    """AIチェック JSON が生成され、期待構造を持つ。"""
    import ai_check_extract

    out = tmp_path / "out.json"
    ai_check_extract.write_ai_check_json(str(sample_pptx), str(out))
    data = json.loads(out.read_text(encoding="utf-8"))

    assert data["format"] == "pptx"
    assert isinstance(data.get("slides"), list)
    assert len(data["slides"]) > 0
    assert "table_of_contents" in data


def test_prompt_template_is_ai_check(tmp_path):
    """プロンプト手順書に AIチェック 用語が含まれる。"""
    import ai_check_extract

    p = tmp_path / "prompt.md"
    ai_check_extract.write_prompt_hint(str(p))
    text = p.read_text(encoding="utf-8")
    assert "AIチェック" in text
    assert "Phase 3" not in text
    assert "_aicheck" in text
    assert "_phase3" not in text


def test_strategy_checks_runs(sample_pptx):
    """strategy_checks が Finding を返す。"""
    import extractors
    import strategy_checks

    doc = extractors.extract(str(sample_pptx))
    findings = strategy_checks.run_strategy_checks(doc)
    # 少なくとも何らかの指摘（空でも型が正しい）
    assert isinstance(findings, list)
    for f in findings:
        assert hasattr(f, "severity")
        assert hasattr(f, "category")
        assert f.category.startswith("strategy/")


def test_strategy_new_checks_importable():
    """追加3関数が export されている。"""
    import strategy_checks
    assert callable(strategy_checks.check_katakana_overload)
    assert callable(strategy_checks.check_assumptions_missing)
    assert callable(strategy_checks.check_executive_summary_missing)


def test_katakana_overload_detects_density(tmp_path):
    """カタカナ密度が高いスライドを拾う。"""
    from pptx import Presentation
    from pptx.util import Inches
    import extractors
    import strategy_checks

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(4))
    tb.text_frame.text = (
        "エンゲージメント、アラインメント、コミットメント、"
        "シナジー、ステークホルダー、イニシアチブ、フレームワーク、"
        "オペレーショナルエクセレンス、ディスラプション"
    )
    p = tmp_path / "k.pptx"
    prs.save(str(p))
    doc = extractors.extract(str(p))
    findings = strategy_checks.check_katakana_overload(doc)
    assert any(f.category == "strategy/katakana-overload" for f in findings)


def test_assumptions_missing_detects(tmp_path):
    """市場規模を予測しているが前提記載がない → 検出。"""
    from pptx import Presentation
    from pptx.util import Inches
    import extractors
    import strategy_checks

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
    tb.text_frame.text = "市場規模は10兆円に拡大する予測。成長率は年率15%"
    p = tmp_path / "a.pptx"
    prs.save(str(p))
    doc = extractors.extract(str(p))
    findings = strategy_checks.check_assumptions_missing(doc)
    assert any(f.category == "strategy/assumptions-missing" for f in findings)


def test_assumptions_present_not_flagged(tmp_path):
    """予測に前提が明記されていれば検出しない。"""
    from pptx import Presentation
    from pptx.util import Inches
    import extractors
    import strategy_checks

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
    tb.text_frame.text = "市場規模は10兆円に拡大する予測 (前提: CAGR 15%, 試算根拠は別紙参照)"
    p = tmp_path / "a2.pptx"
    prs.save(str(p))
    doc = extractors.extract(str(p))
    findings = strategy_checks.check_assumptions_missing(doc)
    assert not any(f.category == "strategy/assumptions-missing" for f in findings)


def test_executive_summary_missing_on_large_deck(tmp_path):
    """5枚以上あるのに冒頭にサマリーがない → 検出。"""
    from pptx import Presentation
    import extractors
    import strategy_checks

    prs = Presentation()
    titles = ["表紙", "目次", "会社概要", "現状分析", "課題", "提案内容", "体制"]
    for t in titles:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        if s.shapes.title:
            s.shapes.title.text = t
    p = tmp_path / "e.pptx"
    prs.save(str(p))
    doc = extractors.extract(str(p))
    findings = strategy_checks.check_executive_summary_missing(doc)
    assert any(f.category == "strategy/executive-summary-missing" for f in findings)


def test_executive_summary_present_not_flagged(tmp_path):
    """冒頭にエグゼクティブサマリーがあれば検出しない。"""
    from pptx import Presentation
    import extractors
    import strategy_checks

    prs = Presentation()
    titles = ["表紙", "エグゼクティブサマリー", "現状", "提案", "体制", "費用"]
    for t in titles:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        if s.shapes.title:
            s.shapes.title.text = t
    p = tmp_path / "e2.pptx"
    prs.save(str(p))
    doc = extractors.extract(str(p))
    findings = strategy_checks.check_executive_summary_missing(doc)
    assert not any(f.category == "strategy/executive-summary-missing" for f in findings)


def test_llm_system_prompt_has_new_axes():
    """LLM プロンプトに追加3観点が入っている。"""
    import llm_review
    p = llm_review.SYSTEM_PROMPT
    assert "ロジックの飛躍" in p
    assert "フレームワーク整合" in p
    assert "実行可能性" in p or "Feasibility" in p
    # カテゴリ識別子も追加されている
    assert "logic-leap" in p
    assert "framework" in p
    assert "feasibility" in p


def test_skill_md_declares_no_external_llm_for_claude_code():
    """SKILL.md に「Claude Code から呼ばれた場合は外部LLM APIを使わない」旨が書かれている。"""
    skill_md = (SKILL_ROOT / "SKILL.md").read_text(encoding="utf-8")
    # workflow 節に Claude 自身がレビューする旨
    assert "Claude 自身" in skill_md or "Claude自身" in skill_md
    # 外部API不使用の明示
    assert "APIキーは不要" in skill_md or "外部API" in skill_md
    # llm_review.py を使わないことの明示
    assert "llm_review.py" in skill_md


def test_cli_guidance_message_mentions_claude_self_review(sample_pptx, tmp_path):
    """review.py --ai-check-json の出力メッセージに Claude 自身がレビューする旨が含まれる。"""
    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "review.py"),
            str(sample_pptx),
            "--skip-liveness",
            "--ai-check-json",
            "--out-dir",
            str(tmp_path),
        ],
        capture_output=True,
        text=True,
        timeout=60,
    )
    assert result.returncode == 0
    combined = result.stdout + result.stderr
    assert "Claude" in combined
    assert "外部LLM API" in combined or "外部API" in combined


def test_cli_ai_check_json_flag(sample_pptx, tmp_path):
    """review.py CLI の --ai-check-json が動く（--skip-liveness で外部通信なし）。"""
    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "review.py"),
            str(sample_pptx),
            "--skip-liveness",
            "--ai-check-json",
            "--out-dir",
            str(tmp_path),
        ],
        capture_output=True,
        text=True,
        timeout=60,
    )
    assert result.returncode == 0, f"CLI failed: {result.stderr}"
    stem = sample_pptx.stem
    assert (tmp_path / f"{stem}_review.md").exists()
    assert (tmp_path / f"{stem}_aicheck.json").exists()
    assert (tmp_path / f"{stem}_aicheck_prompt.md").exists()


def test_cli_old_phase3_flag_rejected(sample_pptx, tmp_path):
    """旧フラグ --phase3-json は受け付けない（リネーム確認）。"""
    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "review.py"),
            str(sample_pptx),
            "--skip-liveness",
            "--phase3-json",
            "--out-dir",
            str(tmp_path),
        ],
        capture_output=True,
        text=True,
        timeout=30,
    )
    assert result.returncode != 0
    assert "unrecognized" in result.stderr.lower() or "error" in result.stderr.lower()


def test_llm_review_missing_api_key(tmp_path, monkeypatch):
    """API キーなしだと丁寧なエラーを返す（APIは叩かない）。"""
    import llm_review

    # 環境変数を空に（実際の .env 探索もブロック）
    monkeypatch.delenv("GOOGLE_API_KEY", raising=False)
    monkeypatch.delenv("GEMINI_API_KEY", raising=False)
    monkeypatch.delenv("DELIVERABLE_REVIEW_ENV_FILE", raising=False)
    # _candidate_env_paths を tmp 内の存在しないパスだけに差し替え
    from pathlib import Path as _Path
    monkeypatch.setattr(
        llm_review,
        "_candidate_env_paths",
        lambda explicit=None: [_Path(str(tmp_path / "nonexistent.env"))],
    )

    # 空の JSON を作成
    j = tmp_path / "empty.json"
    j.write_text("{}", encoding="utf-8")

    findings, err = llm_review.run_llm_review(str(j))
    assert findings == []
    assert err is not None
    assert "API" in err or "KEY" in err.upper()
