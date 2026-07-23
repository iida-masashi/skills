"""Strategy-consulting quality checks (rule-based, local).

戦略コンサルタント品質の機械判定チェック群。
MECE/ピラミッド原則/So What? の厳密判定はLLM（llm_review.py）に任せ、
ここでは正規表現・カウントで判定できる周辺ルールのみ扱う。

Findings は checker="consulting-layout" で発行し、
category="strategy/*" で分類する（既存パイプラインとの統合容易性のため）。
"""
from __future__ import annotations

import re
from collections import Counter
from typing import List

from checkers import Finding, SEVERITY_MEDIUM, SEVERITY_LOW, SEVERITY_INFO
from extractors import iter_shapes_recursive


CHECKER = "consulting-layout"

# タイトル長の上限（目安: 1行に収まる範囲）
TITLE_MAX_CHARS = 40

# 1スライドのバレット/段落数の上限（Miller's Law: 7±2）
BULLETS_MAX = 9

# 結論/提言キーワード
CONCLUSION_WORDS = re.compile(
    r"(まとめ|結論|提言|提案|recommendation|summary|conclusion|next\s*steps|今後|アクション|action\s*plan)",
    re.IGNORECASE,
)

# アジェンダ/目次キーワード
AGENDA_WORDS = re.compile(
    r"(アジェンダ|目次|agenda|contents|本日の内容|table\s+of\s+contents)",
    re.IGNORECASE,
)

# 主語なし・曖昧タイトルの典型
VAGUE_TITLE_PATTERNS = [
    re.compile(r"^.{0,15}について$"),
    re.compile(r"^.{0,15}に関して$"),
    re.compile(r"^(概要|状況|現状|課題|検討|考察|分析|まとめ)$"),
    re.compile(r"^(overview|status|issues?|analysis|considerations?)$", re.IGNORECASE),
]

# セクション見出しに典型的な語
SECTION_LIKE = re.compile(
    r"^\s*(第[一二三四五六七八九十\d]+[章節部]|Chapter\s+\d+|Section\s+\d+|Part\s+[IVX\d]+|\d+[\.\)]\s*)",
    re.IGNORECASE,
)


def _pptx_iter_slides(doc):
    """PPTXの各スライドをyield: (idx, slide, title_text)"""
    if doc.ext != ".pptx" or doc.raw is None:
        return
    from pptx.util import Emu  # noqa: F401
    prs = doc.raw
    for idx, slide in enumerate(prs.slides, start=1):
        title_text = ""
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = (slide.shapes.title.text_frame.text or "").strip()
        except Exception:
            pass
        yield idx, slide, title_text


def check_title_length(doc) -> List[Finding]:
    """タイトルが長すぎる（1スライド1メッセージ違反の兆候）。"""
    out: List[Finding] = []
    for idx, _slide, title in _pptx_iter_slides(doc):
        if not title:
            continue
        if len(title) > TITLE_MAX_CHARS:
            out.append(Finding(
                checker=CHECKER,
                severity=SEVERITY_LOW,
                category="strategy/title-too-long",
                location_label=f"Slide {idx}",
                location_index=idx,
                evidence=title[:80],
                note=f"タイトルが{len(title)}文字。{TITLE_MAX_CHARS}文字以下にメッセージを凝縮することを推奨。",
                source_handle=None,
            ))
    return out


def check_title_multiple_sentences(doc) -> List[Finding]:
    """タイトル内に句点「。」や複数文が混在（1スライド1メッセージ違反）。"""
    out: List[Finding] = []
    for idx, _slide, title in _pptx_iter_slides(doc):
        if not title:
            continue
        # 末尾の句点は除いてカウント
        trimmed = title.rstrip("。.")
        if "。" in trimmed or re.search(r"[。][^」』)\)]", title):
            out.append(Finding(
                checker=CHECKER,
                severity=SEVERITY_LOW,
                category="strategy/title-multi-sentence",
                location_label=f"Slide {idx}",
                location_index=idx,
                evidence=title[:80],
                note="タイトルに句点が複数。1スライド1メッセージを意識し、1文に凝縮することを推奨。",
                source_handle=None,
            ))
    return out


def check_vague_title(doc) -> List[Finding]:
    """曖昧・主語なしタイトルの検出（体言止めチェックとは別観点）。"""
    out: List[Finding] = []
    for idx, _slide, title in _pptx_iter_slides(doc):
        if not title:
            continue
        for pat in VAGUE_TITLE_PATTERNS:
            if pat.match(title.strip()):
                out.append(Finding(
                    checker=CHECKER,
                    severity=SEVERITY_MEDIUM,
                    category="strategy/vague-title",
                    location_label=f"Slide {idx}",
                    location_index=idx,
                    evidence=title[:80],
                    note="主語・結論が不明瞭なタイトル。「何について・どうなのか」を含めるのが戦略コンサル流。",
                    source_handle=None,
                ))
                break
    return out


def check_bullet_overload(doc) -> List[Finding]:
    """1スライド内のバレット/段落が多すぎる（Miller's Law: 7±2超過）。"""
    if doc.ext != ".pptx" or doc.raw is None:
        return []
    out: List[Finding] = []
    for idx, slide, _title in _pptx_iter_slides(doc):
        bullet_count = 0
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None
        for shape in iter_shapes_recursive(slide.shapes):
            if not shape.has_text_frame:
                continue
            if title_shape is not None and shape is title_shape:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                if text:
                    bullet_count += 1
        if bullet_count > BULLETS_MAX:
            out.append(Finding(
                checker=CHECKER,
                severity=SEVERITY_LOW,
                category="strategy/bullet-overload",
                location_label=f"Slide {idx}",
                location_index=idx,
                evidence=f"{bullet_count}行",
                note=f"テキスト行が{bullet_count}行（目安: {BULLETS_MAX}行以下）。情報過多。",
                source_handle=None,
            ))
    return out


def check_conclusion_missing(doc) -> List[Finding]:
    """ファイル全体に結論/提言スライドが1枚もない。"""
    if doc.ext not in (".pptx", ".docx", ".pdf"):
        return []
    has_conclusion = False
    # タイトル＋本文どちらでも判定
    for idx, _slide, title in _pptx_iter_slides(doc):
        if title and CONCLUSION_WORDS.search(title):
            has_conclusion = True
            break
    if not has_conclusion:
        # 本文にも探す
        for text in doc.location_text.values():
            if CONCLUSION_WORDS.search(text or ""):
                has_conclusion = True
                break
    if not has_conclusion and doc.location_flags:
        # 十分な規模（3ページ以上）のときだけ指摘
        if len(doc.location_flags) >= 3:
            return [Finding(
                checker=CHECKER,
                severity=SEVERITY_MEDIUM,
                category="strategy/missing-conclusion",
                location_label="(全体)",
                location_index=0,
                evidence="",
                note="結論/提言/まとめ/Next Steps を示すスライドが見当たらない。戦略提案には必須。",
                source_handle=None,
            )]
    return []


def check_agenda_vs_sections(doc) -> List[Finding]:
    """アジェンダ記載項目とセクション構成の不一致（PPTX限定、簡易）。"""
    if doc.ext != ".pptx" or doc.raw is None:
        return []
    # 1. アジェンダスライドを見つける（タイトル or 本文にagenda/目次があるスライド）
    agenda_idx = None
    agenda_items: list[str] = []
    for idx, slide, title in _pptx_iter_slides(doc):
        body_text = ""
        for shape in iter_shapes_recursive(slide.shapes):
            if not shape.has_text_frame:
                continue
            body_text += "\n" + (shape.text_frame.text or "")
        combined = (title or "") + "\n" + body_text
        if AGENDA_WORDS.search(combined):
            agenda_idx = idx
            try:
                title_shape = slide.shapes.title
            except Exception:
                title_shape = None
            for shape in iter_shapes_recursive(slide.shapes):
                if not shape.has_text_frame:
                    continue
                if title_shape is not None and shape is title_shape:
                    continue
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    # 番号/記号先頭の除去
                    t = re.sub(r"^[\s\d\.\)・\-\*●○◆◇■□•]+", "", t).strip()
                    if len(t) >= 3:
                        agenda_items.append(t)
            break
    if agenda_idx is None or len(agenda_items) < 2:
        return []
    # 2. 後続スライドのタイトル群とマッチングを試みる（部分一致）
    titles_after: list[str] = []
    for idx, _slide, title in _pptx_iter_slides(doc):
        if idx <= agenda_idx or not title:
            continue
        titles_after.append(title)
    if not titles_after:
        return []
    # アジェンダ各項目のうち、後続タイトルに部分一致で現れないものをカウント
    missing = []
    for item in agenda_items:
        # 短いトークンに分解（2文字以上）
        tokens = [w for w in re.split(r"[\s、,・／/]+", item) if len(w) >= 2]
        # 代表語としてアイテム全体 or 最長トークンで検索
        search_key = item if len(item) <= 10 else (max(tokens, key=len) if tokens else item)
        found = any(search_key in t for t in titles_after)
        if not found:
            missing.append(item)
    if len(missing) >= 2:  # 2件以上マッチしなければ警告
        return [Finding(
            checker=CHECKER,
            severity=SEVERITY_LOW,
            category="strategy/agenda-mismatch",
            location_label=f"Slide {agenda_idx}",
            location_index=agenda_idx,
            evidence=" / ".join(missing[:3]),
            note=f"アジェンダ項目のうち{len(missing)}件が後続スライドのタイトルに見当たらない。構成とアジェンダの整合性を確認。",
            source_handle=None,
        )]
    return []


def check_deck_too_sparse(doc) -> List[Finding]:
    """スライドが1枚しかない等、資料として成立しているか簡易チェック。"""
    if doc.ext != ".pptx":
        return []
    n = len(doc.location_flags)
    if 2 <= n <= 3:
        return [Finding(
            checker=CHECKER,
            severity=SEVERITY_INFO,
            category="strategy/deck-too-sparse",
            location_label="(全体)",
            location_index=0,
            evidence=f"{n}スライド",
            note=f"スライド{n}枚のみ。構成が「現状→課題→解決→効果」で成立しているか確認。",
            source_handle=None,
        )]
    return []


# ----------------------------------------------------------------------------
# 追加チェック (2026-04 実装)
# ----------------------------------------------------------------------------

# カタカナ多用判定の閾値
KATAKANA_DENSITY_THRESHOLD = 0.20   # 1スライド全テキストの20%以上がカタカナ
KATAKANA_MIN_CHARS = 40             # ノイズ抑制: 全テキスト40文字未満はスキップ
KATAKANA_UNIQUE_WORDS_MIN = 5       # ユニークなカタカナ語が5種以上

# カタカナ連続ブロック（2文字以上）
_KATAKANA_BLOCK = re.compile(r"[゠-ヿ]{2,}")
_KATAKANA_CHAR = re.compile(r"[゠-ヿ]")

# 許容リスト（固有名詞的なカタカナ語は除外）— 拡張可
KATAKANA_ALLOWLIST = {
    "ソリューション", "システム", "サービス", "データ", "プロジェクト",
    "ビジネス", "クライアント", "パートナー",
}

# 前提条件キーワード
ASSUMPTION_WORDS = re.compile(
    r"(前提|仮定|仮説|試算|概算|推計|assumption|hypothesis|estimate|"
    r"見込み|想定|シナリオ|scenario|前年比|対前年|YoY|CAGR)",
    re.IGNORECASE,
)

# 予測・見通しを含むキーワード（前提条件が必要なコンテキストを検出）
FORECAST_WORDS = re.compile(
    r"(予測|予想|見通し|将来|forecast|projection|\d{4}年度?予測|"
    r"\d+年後|\d+%\s*(成長|増加|拡大)|伸び率|市場規模|TAM|SAM|SOM)",
    re.IGNORECASE,
)

# エグゼクティブサマリー判定キーワード
EXEC_SUMMARY_WORDS = re.compile(
    r"(エグゼクティブ\s*サマリー|executive\s*summary|要旨|要約|"
    r"本提案の要点|本日の結論|本資料のポイント|key\s*messages?|"
    r"主要メッセージ|3つの(ポイント|メッセージ|キー))",
    re.IGNORECASE,
)

# カバー/目次/会社概要 と判別するためのキーワード（サマリーの位置判定に使用）
COVER_OR_AGENDA = re.compile(
    r"(表紙|cover|目次|アジェンダ|agenda|contents|会社概要|company\s*profile)",
    re.IGNORECASE,
)


def _extract_all_text_pptx(slide) -> str:
    """PPTXのスライド全テキストを連結（タイトル含む）。"""
    buf = []
    for shape in iter_shapes_recursive(slide.shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text or ""
        if t:
            buf.append(t)
    return "\n".join(buf)


def check_katakana_overload(doc) -> List[Finding]:
    """カタカナ多用スライドの検出。"""
    if doc.ext != ".pptx":
        return []
    out: List[Finding] = []
    for idx, slide, _title in _pptx_iter_slides(doc):
        text = _extract_all_text_pptx(slide)
        if len(text) < KATAKANA_MIN_CHARS:
            continue
        # カタカナ文字数密度
        kata_chars = len(_KATAKANA_CHAR.findall(text))
        density = kata_chars / max(1, len(text))
        # カタカナ語のユニーク数（許容リスト除外）
        blocks = _KATAKANA_BLOCK.findall(text)
        unique_words = {b for b in blocks if b not in KATAKANA_ALLOWLIST and len(b) >= 3}
        if density >= KATAKANA_DENSITY_THRESHOLD and len(unique_words) >= KATAKANA_UNIQUE_WORDS_MIN:
            top = sorted(unique_words, key=len, reverse=True)[:5]
            out.append(Finding(
                checker=CHECKER,
                severity=SEVERITY_LOW,
                category="strategy/katakana-overload",
                location_label=f"Slide {idx}",
                location_index=idx,
                evidence=" / ".join(top),
                note=f"カタカナ密度{density:.0%}・{len(unique_words)}語。"
                     "日本語で言い換え可能な用語は置換し、稚拙な印象を回避。",
                source_handle=None,
            ))
    return out


def check_assumptions_missing(doc) -> List[Finding]:
    """予測・試算を含むスライドで前提条件の明示がない。"""
    if doc.ext != ".pptx":
        return []
    out: List[Finding] = []
    for idx, slide, title in _pptx_iter_slides(doc):
        text = (title or "") + "\n" + _extract_all_text_pptx(slide)
        if not FORECAST_WORDS.search(text):
            continue
        # スピーカーノートも検索対象に含める（補足に書いてあれば許容）
        try:
            if slide.has_notes_slide:
                text += "\n" + (slide.notes_slide.notes_text_frame.text or "")
        except Exception:
            pass
        if ASSUMPTION_WORDS.search(text):
            continue
        out.append(Finding(
            checker=CHECKER,
            severity=SEVERITY_MEDIUM,
            category="strategy/assumptions-missing",
            location_label=f"Slide {idx}",
            location_index=idx,
            evidence=(title or "")[:60],
            note="予測・見通し・市場規模等を述べているが、前提/仮定/算出根拠の記載がない。"
                 "「前提: 〜」「試算根拠: 〜」を明記すること。",
            source_handle=None,
        ))
    return out


def check_executive_summary_missing(doc) -> List[Finding]:
    """エグゼクティブサマリーが冒頭にない（PPTXで5枚以上の資料限定）。"""
    if doc.ext != ".pptx":
        return []
    total = len(doc.location_flags)
    if total < 5:
        return []  # 小規模資料は対象外
    # Slide 1-4 の範囲にサマリー的タイトル or 本文があるかチェック
    has_summary = False
    for idx, slide, title in _pptx_iter_slides(doc):
        if idx > 4:
            break
        text = (title or "") + "\n" + _extract_all_text_pptx(slide)
        # cover/agenda/会社概要はサマリーとして扱わない
        if EXEC_SUMMARY_WORDS.search(text) and not (
            title and COVER_OR_AGENDA.search(title)
        ):
            has_summary = True
            break
    if not has_summary:
        return [Finding(
            checker=CHECKER,
            severity=SEVERITY_MEDIUM,
            category="strategy/executive-summary-missing",
            location_label="(全体)",
            location_index=0,
            evidence=f"{total}スライド",
            note="冒頭4スライド以内にエグゼクティブサマリー/要旨/キーメッセージが見当たらない。"
                 "経営層が30秒で判断できるよう、結論先行(PREP)で冒頭に配置すべき。",
            source_handle=None,
        )]
    return []


def run_strategy_checks(doc) -> List[Finding]:
    """全ての戦略品質チェックを実行。"""
    findings: List[Finding] = []
    findings.extend(check_title_length(doc))
    findings.extend(check_title_multiple_sentences(doc))
    findings.extend(check_vague_title(doc))
    findings.extend(check_bullet_overload(doc))
    findings.extend(check_conclusion_missing(doc))
    findings.extend(check_agenda_vs_sections(doc))
    findings.extend(check_deck_too_sparse(doc))
    findings.extend(check_katakana_overload(doc))
    findings.extend(check_assumptions_missing(doc))
    findings.extend(check_executive_summary_missing(doc))
    return findings
