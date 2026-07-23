"""AIチェック: structured JSON extraction of slide/page content for an LLM
(Claude Code / Gemini / etc.) to perform qualitative consulting review —
pyramid principle, MECE, So What? / Why So?, action concreteness, narrative flow.

This module does NOT call any LLM. It produces a structured JSON that a
reviewing agent can read and reason over.
"""
import json
import re
from typing import Dict, Any, List
from pathlib import Path

import extractors
from extractors import iter_shapes_recursive


def _classify_slide_role(title: str, body: str) -> str:
    """Rough heuristic role classification."""
    t = (title or "").lower()
    b = (body or "")[:200]
    if any(k in title for k in ("目次", "アジェンダ", "Agenda", "Contents")):
        return "agenda"
    if any(k in title for k in ("会社概要", "Company", "企業概要")):
        return "company-profile"
    if any(k in title for k in ("経歴", "Profile", "プロフィール")):
        return "team-profile"
    if any(k in title for k in ("まとめ", "Summary", "Conclusion")):
        return "conclusion"
    if any(k in title for k in ("Appendix", "参考", "補足", "付録")):
        return "appendix"
    if any(k in title for k in ("表紙", "Cover")):
        return "cover"
    if any(k in title for k in ("課題", "問題", "現状")):
        return "problem"
    if any(k in title for k in ("提案", "解決", "アプローチ", "方針", "対応")):
        return "solution"
    if any(k in title for k in ("効果", "成果", "ベネフィット", "Benefit", "Impact")):
        return "benefit"
    if any(k in title for k in ("スケジュール", "計画", "Plan", "Roadmap")):
        return "schedule"
    if any(k in title for k in ("体制", "チーム", "Team")):
        return "team"
    if any(k in title for k in ("費用", "金額", "見積", "Cost", "Fee")):
        return "pricing"
    return "content"


def extract_structure_pptx(path: str) -> Dict[str, Any]:
    from pptx import Presentation
    prs = Presentation(path)
    slides_json = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        try:
            t = slide.shapes.title
            if t is not None:
                title = (t.text or "").strip()
        except Exception:
            pass

        body_paragraphs: List[str] = []
        bullets: List[str] = []
        tables: List[List[List[str]]] = []
        has_image = False
        has_chart = False
        has_table = False
        is_hidden = slide.element.get("show") == "0"

        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            pass
        for shape in iter_shapes_recursive(slide.shapes):
            if title_shape is not None and shape is title_shape:
                continue
            try:
                if shape.shape_type == 13:
                    has_image = True
            except Exception:
                pass
            if shape.has_chart:
                has_chart = True
            if shape.has_table:
                has_table = True
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text for cell in row.cells])
                tables.append(rows)
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs) or ""
                    t = t.strip()
                    if not t:
                        continue
                    body_paragraphs.append(t)
                    try:
                        if para.level and para.level > 0:
                            bullets.append(t)
                    except Exception:
                        pass

        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()

        body_text = "\n".join(body_paragraphs)
        slides_json.append({
            "slide": idx,
            "hidden": is_hidden,
            "title": title,
            "role": _classify_slide_role(title, body_text),
            "body_paragraphs": body_paragraphs,
            "bullets": bullets,
            "tables": tables,
            "has_image": has_image,
            "has_chart": has_chart,
            "has_table": has_table,
            "speaker_notes": notes,
        })

    toc = [{"slide": s["slide"], "title": s["title"], "role": s["role"]}
           for s in slides_json if not s["hidden"]]

    return {
        "format": "pptx",
        "path": path,
        "slide_count": len(slides_json),
        "table_of_contents": toc,
        "slides": slides_json,
    }


def extract_structure_docx(path: str) -> Dict[str, Any]:
    from docx import Document as DocxDocument
    d = DocxDocument(path)
    paragraphs = []
    for p in d.paragraphs:
        t = (p.text or "").strip()
        if not t:
            continue
        style = (p.style.name or "") if p.style else ""
        is_heading = style.startswith("Heading") or style.startswith("見出し")
        paragraphs.append({
            "text": t,
            "style": style,
            "is_heading": is_heading,
            "level": int(re.search(r"\d+", style).group(0)) if is_heading and re.search(r"\d+", style) else 0,
        })
    tables = []
    for tbl in d.tables:
        rows = [[cell.text for cell in row.cells] for row in tbl.rows]
        tables.append(rows)
    return {
        "format": "docx",
        "path": path,
        "paragraphs": paragraphs,
        "tables": tables,
    }


def extract_structure_pdf(path: str) -> Dict[str, Any]:
    import pdfplumber
    pages = []
    with pdfplumber.open(path) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            tables = []
            try:
                tables = page.extract_tables() or []
            except Exception:
                pass
            pages.append({
                "page": idx,
                "text": text,
                "tables": tables,
            })
    return {
        "format": "pdf",
        "path": path,
        "page_count": len(pages),
        "pages": pages,
    }


def extract_structure(path: str) -> Dict[str, Any]:
    lower = path.lower()
    if lower.endswith(".pptx"):
        return extract_structure_pptx(path)
    if lower.endswith(".docx"):
        return extract_structure_docx(path)
    if lower.endswith(".pdf"):
        return extract_structure_pdf(path)
    raise ValueError(f"Unsupported file type: {path}")


def write_ai_check_json(path: str, out_path: str) -> None:
    data = extract_structure(path)
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def _load_strategic_review_prompt() -> str:
    """戦略レビュー用 SYSTEM_PROMPT を llm_review から取り込む（単一ソース）。

    llm_review.py がインポートに失敗する場合（依存欠落等）はフォールバック
    テキストを返す。
    """
    try:
        from llm_review import SYSTEM_PROMPT
        return SYSTEM_PROMPT
    except Exception:
        return (
            "（llm_review.SYSTEM_PROMPT のインポートに失敗しました。"
            "scripts/llm_review.py の SYSTEM_PROMPT を直接参照してください。）"
        )


_CLAUDE_CODE_ADDENDUM = """\

# Claude Code 向け追加指示（このスキルから呼ばれた場合）

あなたは外部 LLM API ではなく Claude Code 内蔵のレビュアーとして動作しています。
以下のルールに従ってください:

## 入力
同フォルダ内の `*_aicheck.json` を `Read` ツールで読み込み、`slides[*]` と
`table_of_contents` を全量見ること。bullets/body/tables/speaker_notes を全部
読む（要約しない）。

## 出力
レビュー結果は `*_aicheck_review.md` として書き出すこと。フォーマットは
SYSTEM_PROMPT の JSON ではなく、人間が読みやすい Markdown:

```markdown
# 戦略レビュー（Tier-1 MD/Partner視点）

## 総評
- **評価**: A / B / C / D
- **提出可否**: 提出可 / 要修正 / 大幅手直し必要
- **Key Question**: <この資料が答えようとしている中心問い>
- **答え明瞭度**: 明確 / 推測可能 / 不明
- **ストーリー要旨**: <冒頭〜結論を3〜5文に圧縮>
- **強み (Top 2)**: ...
- **弱み (Top 3)**: ...
- **Partner一言**: <現場感ある1行コメント>

## 観点別レビュー（15観点・全て言及）

### 1. ピラミッド原則 (pyramid)
- **Severity**: HIGH / MEDIUM / LOW / INFO
- **Slide**: 該当スライド番号
- **引用**: 「<原文を引用>」
- **指摘**: <何が問題か。断定形>
- **影響**: <クライアントに渡るとどう不利益か>
- **改善案**: <どう直すか>
- **書き換え例**: 現状「X」→ 改善「Y」

（以下、mece / so-what / issue-tree / logic-leap / data-rigor / framework /
action / feasibility / balance / client-view / alternatives / premise /
story-line / risk-scenario の14観点を同形式で。問題のない観点も
「該当なし: <根拠>」を必ず1行は書く。観点を黙ってスキップしないこと。）

## 優先アクション (Top 5)
1. 〔Slide N〕 <最も重要な修正>
2. ...
```

## 品質基準
- 指摘は断定形で書く。「〜の可能性がある」「〜と思われる」は使わない。
- 各指摘には**スライド番号と原文引用**を必ず添える。
- 書き換え例（現状→改善）をタイトル/メッセージ系の指摘では必ず示す。
- Severity は SYSTEM_PROMPT のルーブリックに従って機械的に決める。

## 禁止事項
- 外部APIを呼ばない（`llm_review.py` を実行しない）。
- JSON 内の情報以外は使わない。推測の域を出ない場合は「判断不能」と明記。
"""


AI_CHECK_PROMPT_TEMPLATE = (
    "# AIチェック 定性レビュー指示（戦略コンサル品質）\n\n"
    "この `*_aicheck.json` は、deliverable-review スキルが抽出した資料の構造データです。\n"
    "Tier-1 戦略コンサルの MD/Partner として、以下の SYSTEM_PROMPT に従い"
    "**赤入れレベルの厳しさ**で 15 観点の定性レビューを行ってください。\n\n"
    "---\n\n"
    "# SYSTEM_PROMPT（llm_review.py と完全同一・単一ソース）\n\n"
    + _load_strategic_review_prompt()
    + _CLAUDE_CODE_ADDENDUM
)


def write_prompt_hint(out_path: str) -> None:
    # 動的に最新の SYSTEM_PROMPT を取り込んで書き出す（モジュール import 時の
    # キャッシュではなく、書き出すたびに最新化する）
    template = (
        "# AIチェック 定性レビュー指示（戦略コンサル品質）\n\n"
        "この `*_aicheck.json` は、deliverable-review スキルが抽出した資料の構造データです。\n"
        "Tier-1 戦略コンサルの MD/Partner として、以下の SYSTEM_PROMPT に従い"
        "**赤入れレベルの厳しさ**で 15 観点の定性レビューを行ってください。\n\n"
        "---\n\n"
        "# SYSTEM_PROMPT（llm_review.py と完全同一・単一ソース）\n\n"
        + _load_strategic_review_prompt()
        + _CLAUDE_CODE_ADDENDUM
    )
    Path(out_path).write_text(template, encoding="utf-8")
