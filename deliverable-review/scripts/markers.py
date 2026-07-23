"""Write review-copy with visual markers to .pptx / .docx.

PDF is not supported for marking (report-only per user request).
"""
import shutil
from collections import defaultdict
from typing import List


SEVERITY_COLOR = {
    "HIGH":   (0xE6, 0x00, 0x23),   # red
    "MEDIUM": (0xFF, 0x8C, 0x00),   # orange
    "LOW":    (0xFF, 0xD7, 0x00),   # yellow
    "INFO":   (0x4F, 0x81, 0xBD),   # blue
}

SEVERITY_ORDER = {"HIGH": 0, "MEDIUM": 1, "LOW": 2, "INFO": 3}


def _worst_severity(findings):
    return min(findings, key=lambda f: SEVERITY_ORDER.get(f.severity, 99)).severity


# ------------------------------------------------------------
# PPTX marker
# ------------------------------------------------------------

def mark_pptx(src_path: str, dst_path: str, findings: List) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    shutil.copyfile(src_path, dst_path)
    prs = Presentation(dst_path)

    # Group by slide index
    by_slide = defaultdict(list)
    for f in findings:
        by_slide[f.location_index].append(f)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_findings = by_slide.get(slide_idx, [])
        if not slide_findings:
            continue

        worst = _worst_severity(slide_findings)
        r, g, b = SEVERITY_COLOR[worst]

        # Add an annotation text box in the top-right corner
        slide_w = prs.slide_width
        box_w = Inches(3.2)
        box_h = Inches(0.4 + 0.18 * min(len(slide_findings), 10))
        left = slide_w - box_w - Emu(91440)  # 0.1 inch margin
        top = Emu(91440)

        tb = slide.shapes.add_textbox(left, top, box_w, box_h)
        tf = tb.text_frame
        tf.word_wrap = True

        header = tf.paragraphs[0]
        hr = header.add_run()
        hr.text = f"⚠ REVIEW [{worst}] ({len(slide_findings)}件)"
        hr.font.bold = True
        hr.font.size = Pt(10)
        hr.font.color.rgb = RGBColor(r, g, b)

        for f in slide_findings[:10]:
            p = tf.add_paragraph()
            run = p.add_run()
            evidence = f.evidence if len(f.evidence) < 60 else f.evidence[:60] + "…"
            run.text = f"• [{f.severity}] {f.checker}/{f.category}: {evidence}"
            run.font.size = Pt(8)
            sr, sg, sb = SEVERITY_COLOR[f.severity]
            run.font.color.rgb = RGBColor(sr, sg, sb)
        if len(slide_findings) > 10:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"…and {len(slide_findings) - 10} more. See report."
            run.font.size = Pt(8)
            run.font.italic = True

        # Colored border on annotation box
        line = tb.line
        line.color.rgb = RGBColor(r, g, b)
        line.width = Pt(1.5)

    prs.save(dst_path)


# ------------------------------------------------------------
# DOCX marker
# ------------------------------------------------------------

def mark_docx(src_path: str, dst_path: str, findings: List) -> None:
    from docx import Document as DocxDocument
    from docx.shared import RGBColor, Pt

    shutil.copyfile(src_path, dst_path)
    d = DocxDocument(dst_path)

    if not findings:
        d.save(dst_path)
        return

    # Insert a summary box at the top of the document
    worst = _worst_severity(findings)
    r, g, b = SEVERITY_COLOR[worst]

    # Build summary paragraphs, then move them to top
    summary_paragraphs = []
    intro = d.add_paragraph()
    run = intro.add_run(f"⚠ REVIEW SUMMARY [{worst}] — {len(findings)}件の指摘")
    run.bold = True
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(r, g, b)
    summary_paragraphs.append(intro)

    # Group and list first 30
    by_sev = defaultdict(list)
    for f in findings:
        by_sev[f.severity].append(f)

    shown = 0
    for sev in ("HIGH", "MEDIUM", "LOW", "INFO"):
        items = by_sev.get(sev, [])
        if not items:
            continue
        header = d.add_paragraph()
        hr, hg, hb = SEVERITY_COLOR[sev]
        hrun = header.add_run(f"[{sev}] {len(items)}件")
        hrun.bold = True
        hrun.font.color.rgb = RGBColor(hr, hg, hb)
        summary_paragraphs.append(header)

        for f in items:
            if shown >= 30:
                break
            p = d.add_paragraph()
            evidence = f.evidence if len(f.evidence) < 80 else f.evidence[:80] + "…"
            prun = p.add_run(f"• {f.location_label} {f.checker}/{f.category}: {evidence}")
            prun.font.size = Pt(9)
            prun.font.color.rgb = RGBColor(hr, hg, hb)
            summary_paragraphs.append(p)
            shown += 1

    if len(findings) > shown:
        more = d.add_paragraph()
        run = more.add_run(f"…and {len(findings) - shown} more. See full report (.md).")
        run.font.size = Pt(9)
        run.italic = True
        summary_paragraphs.append(more)

    sep = d.add_paragraph()
    sep_run = sep.add_run("─" * 40)
    sep_run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    summary_paragraphs.append(sep)

    # Move summary_paragraphs to the beginning of the document body
    body = d.element.body
    for p in reversed(summary_paragraphs):
        body.insert(0, p._element)

    d.save(dst_path)
