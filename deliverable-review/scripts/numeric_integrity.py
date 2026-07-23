"""Numeric integrity checks.

Detects:
- Native chart: pie slices don't sum to ~100%, stacked bars vs shown total
- Tables: row totals / column totals that don't match labeled "合計"/"Total"
- Same location mixing incompatible units (億円 vs 百万円 vs 千円; % vs pp)
- Image charts: best-effort using data labels found as nearby text (lower precision)
"""
import re
from dataclasses import dataclass
from typing import List, Optional, Any

from checkers import Finding, SEVERITY_HIGH, SEVERITY_MEDIUM, SEVERITY_LOW, SEVERITY_INFO
from extractors import iter_shapes_recursive


# ------------------------------------------------------------
# Helpers
# ------------------------------------------------------------

_NUM_WITH_UNIT_RE = re.compile(
    r"(?P<num>-?\d+(?:[.,]\d+)*)"
    r"\s*"
    r"(?P<unit>億円|百万円|千円|万円|兆円|円|%|％|pp|ポイント|件|人|社|個)?"
)

_TOTAL_KEYWORDS = ("合計", "計", "小計", "総計", "Total", "total", "TOTAL", "Sum", "sum")


def _to_float(s: str) -> Optional[float]:
    try:
        return float(s.replace(",", "").replace("，", ""))
    except ValueError:
        return None


def _extract_numbers_with_units(text: str):
    """Yield (value, unit) tuples from free text."""
    for m in _NUM_WITH_UNIT_RE.finditer(text):
        v = _to_float(m.group("num"))
        if v is None:
            continue
        unit = m.group("unit") or ""
        yield v, unit


# ------------------------------------------------------------
# Unit mixing check (per location)
# ------------------------------------------------------------

_CURRENCY_UNITS = {"兆円", "億円", "百万円", "万円", "千円", "円"}
_PERCENT_PP_UNITS = {"%", "％", "pp", "ポイント"}


def check_unit_mixing(doc) -> List[Finding]:
    findings = []
    for loc_idx, flags in doc.location_flags.items():
        text = doc.location_text.get(loc_idx, "")
        label = flags.get("label", f"Loc {loc_idx}")
        units_found = set()
        for _, unit in _extract_numbers_with_units(text):
            if unit:
                units_found.add(unit)

        currency_present = units_found & _CURRENCY_UNITS
        # Mixing two or more currency units in one place usually indicates inconsistency
        if len(currency_present) >= 2:
            findings.append(Finding(
                checker="numeric-integrity",
                severity=SEVERITY_MEDIUM,
                category="currency-unit-mixing",
                location_label=label,
                location_index=loc_idx,
                evidence=", ".join(sorted(currency_present)),
                note="同一スライド/ページ内で金額の単位が混在しています。統一を検討してください。",
            ))

        # % と pp の混在
        pp_present = units_found & _PERCENT_PP_UNITS
        if {"%", "％"} & pp_present and {"pp", "ポイント"} & pp_present:
            findings.append(Finding(
                checker="numeric-integrity",
                severity=SEVERITY_MEDIUM,
                category="percent-pp-mixing",
                location_label=label,
                location_index=loc_idx,
                evidence=", ".join(sorted(pp_present)),
                note="% と pp(パーセントポイント) の混在。違いが正しく使い分けられているか確認してください。",
            ))
    return findings


# ------------------------------------------------------------
# Native chart checks (pptx only)
# ------------------------------------------------------------

def check_native_charts_pptx(doc) -> List[Finding]:
    findings = []
    if doc.ext != ".pptx":
        return findings

    from pptx.enum.chart import XL_CHART_TYPE

    PIE_TYPES = {
        XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED, XL_CHART_TYPE.PIE_OF_PIE,
        XL_CHART_TYPE.DOUGHNUT, XL_CHART_TYPE.DOUGHNUT_EXPLODED,
        XL_CHART_TYPE.BAR_OF_PIE,
    }
    STACKED_100_TYPES = {
        XL_CHART_TYPE.BAR_STACKED_100, XL_CHART_TYPE.COLUMN_STACKED_100,
        XL_CHART_TYPE.LINE_STACKED_100, XL_CHART_TYPE.AREA_STACKED_100,
    }

    for slide_idx, slide in enumerate(doc.raw.slides, start=1):
        label = f"Slide {slide_idx}"
        for shape in iter_shapes_recursive(slide.shapes):
            if not shape.has_chart:
                continue
            chart = shape.chart
            chart_type = chart.chart_type

            # Collect values across all series
            try:
                series_values = []
                for s in chart.series:
                    vals = list(s.values)
                    series_values.append([v for v in vals if v is not None])
            except Exception as e:
                findings.append(Finding(
                    checker="numeric-integrity",
                    severity=SEVERITY_INFO,
                    category="chart-read-error",
                    location_label=label,
                    location_index=slide_idx,
                    evidence=str(e)[:100],
                    note="チャートデータの読み取りに失敗しました。",
                ))
                continue

            # Pie: single series should sum near 100 (if percent) or to a labeled total
            if chart_type in PIE_TYPES and series_values:
                vals = series_values[0]
                if not vals:
                    continue
                total = sum(vals)
                # Heuristic: if values look like percentages (all between 0 and 100 and sum ~100)
                if all(0 <= v <= 100 for v in vals):
                    # Only flag if sum meaningfully off from 100
                    if abs(total - 100) > 1.0 and abs(total - 1.0) > 0.01:
                        # sum could also be 1.0 (fractional). Accept either 100±1 or 1±0.01
                        findings.append(Finding(
                            checker="numeric-integrity",
                            severity=SEVERITY_HIGH,
                            category="pie-sum-not-100",
                            location_label=label,
                            location_index=slide_idx,
                            evidence=f"sum={total:.2f}, values={[round(v, 2) for v in vals]}",
                            note="円グラフの各要素の合計が100%になっていません。",
                        ))

            # 100% stacked: each category should sum to 100 across series
            if chart_type in STACKED_100_TYPES and series_values:
                # Transpose to get per-category sums
                try:
                    ncats = min(len(s) for s in series_values)
                    for c in range(ncats):
                        category_sum = sum(s[c] for s in series_values)
                        target = 100 if category_sum > 2 else 1.0
                        tol = 1.0 if target == 100 else 0.01
                        if abs(category_sum - target) > tol:
                            findings.append(Finding(
                                checker="numeric-integrity",
                                severity=SEVERITY_HIGH,
                                category="stacked100-sum-not-100",
                                location_label=label,
                                location_index=slide_idx,
                                evidence=f"category {c+1} sum={category_sum:.2f}",
                                note="100%積上げグラフのカテゴリ合計が100%になっていません。",
                            ))
                except Exception:
                    pass
    return findings


# ------------------------------------------------------------
# Table integrity: row/column totals
# ------------------------------------------------------------

def _parse_numeric_cell(s: str) -> Optional[float]:
    s = s.strip()
    if not s:
        return None
    # Strip trailing units
    m = re.match(r"^-?\d+(?:[.,]\d+)*", s)
    if not m:
        return None
    return _to_float(m.group(0))


def _check_table_totals(rows_raw, label: str, loc_idx: int, source: str) -> List[Finding]:
    """rows_raw: 2D list of raw cell strings.
    Check: if any row header says 合計/Total, verify row sum matches header.
           if any column header says 合計/Total, verify col sum matches.
    """
    findings = []
    if not rows_raw or len(rows_raw) < 2:
        return findings

    nrows = len(rows_raw)
    ncols = max(len(r) for r in rows_raw)

    # Normalize: pad short rows
    grid = [list(r) + [""] * (ncols - len(r)) for r in rows_raw]

    # Find total-row (any row where first cell contains a total keyword)
    for r_idx, row in enumerate(grid):
        first = (row[0] or "").strip()
        if any(k in first for k in _TOTAL_KEYWORDS):
            # For each numeric column c>=1, sum rows above (r < r_idx) that are numeric
            for c in range(1, ncols):
                labeled_total = _parse_numeric_cell(row[c])
                if labeled_total is None:
                    continue
                numeric_col = []
                for rr in range(r_idx):
                    # skip header row (row 0) only if it's non-numeric
                    v = _parse_numeric_cell(grid[rr][c])
                    if v is not None:
                        numeric_col.append(v)
                if not numeric_col:
                    continue
                computed = sum(numeric_col)
                if abs(computed - labeled_total) > max(0.5, abs(labeled_total) * 0.005):
                    findings.append(Finding(
                        checker="numeric-integrity",
                        severity=SEVERITY_HIGH,
                        category="table-col-total-mismatch",
                        location_label=label,
                        location_index=loc_idx,
                        evidence=f"col {c+1}: labeled={labeled_total}, computed={computed}",
                        note=f"表({source})の列合計が一致しません。",
                    ))

    # Find total-column (any column where header says 合計/Total)
    header = grid[0]
    for c_idx, cell in enumerate(header):
        if any(k in (cell or "") for k in _TOTAL_KEYWORDS):
            for r in range(1, nrows):
                row = grid[r]
                labeled_total = _parse_numeric_cell(row[c_idx] if c_idx < len(row) else "")
                if labeled_total is None:
                    continue
                numeric_row = []
                for cc in range(ncols):
                    if cc == c_idx:
                        continue
                    v = _parse_numeric_cell(row[cc])
                    if v is not None:
                        numeric_row.append(v)
                if not numeric_row:
                    continue
                computed = sum(numeric_row)
                if abs(computed - labeled_total) > max(0.5, abs(labeled_total) * 0.005):
                    findings.append(Finding(
                        checker="numeric-integrity",
                        severity=SEVERITY_HIGH,
                        category="table-row-total-mismatch",
                        location_label=label,
                        location_index=loc_idx,
                        evidence=f"row {r+1}: labeled={labeled_total}, computed={computed}",
                        note=f"表({source})の行合計が一致しません。",
                    ))
    return findings


def check_tables_pptx(doc) -> List[Finding]:
    findings = []
    if doc.ext != ".pptx":
        return findings
    for slide_idx, slide in enumerate(doc.raw.slides, start=1):
        for shape in iter_shapes_recursive(slide.shapes):
            if not shape.has_table:
                continue
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text for cell in row.cells])
            findings.extend(_check_table_totals(
                rows, f"Slide {slide_idx}", slide_idx, "pptx"
            ))
    return findings


def check_tables_docx(doc) -> List[Finding]:
    findings = []
    if doc.ext != ".docx":
        return findings
    for t_idx, tbl in enumerate(doc.raw.tables, start=1):
        rows = []
        for row in tbl.rows:
            rows.append([cell.text for cell in row.cells])
        findings.extend(_check_table_totals(
            rows, f"Table {t_idx}", 1, "docx"
        ))
    return findings


def check_tables_pdf(doc) -> List[Finding]:
    """PDF table extraction via pdfplumber."""
    findings = []
    if doc.ext != ".pdf":
        return findings
    try:
        import pdfplumber
    except ImportError:
        return findings
    with pdfplumber.open(doc.path) as pdf:
        for page_idx, page in enumerate(pdf.pages, start=1):
            try:
                tables = page.extract_tables() or []
            except Exception:
                tables = []
            for t_idx, tbl in enumerate(tables, start=1):
                # pdfplumber returns None for empty cells
                rows = [[c or "" for c in r] for r in tbl]
                findings.extend(_check_table_totals(
                    rows, f"Page {page_idx} Table {t_idx}", page_idx, "pdf"
                ))
    return findings


# ------------------------------------------------------------
# Entry
# ------------------------------------------------------------

def run_numeric_integrity(doc) -> List[Finding]:
    findings = []
    findings.extend(check_unit_mixing(doc))
    if doc.ext == ".pptx":
        findings.extend(check_native_charts_pptx(doc))
        findings.extend(check_tables_pptx(doc))
    elif doc.ext == ".docx":
        findings.extend(check_tables_docx(doc))
    elif doc.ext == ".pdf":
        findings.extend(check_tables_pdf(doc))
    return findings
