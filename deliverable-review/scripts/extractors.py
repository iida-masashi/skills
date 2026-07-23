"""Text extraction from .pptx / .docx / .pdf with location info.

Each extractor yields TextUnit records. A TextUnit is a logical block of text
associated with a location (slide index, paragraph index, or page number) so
later phases can report where findings came from and, for pptx/docx, write
marks back to the originating shape/paragraph.
"""
from dataclasses import dataclass, field
from typing import List, Optional, Any


@dataclass
class TextUnit:
    kind: str                      # "text" | "table-cell" | "image-caption-slot"
    text: str
    location_label: str            # human-readable, e.g. "Slide 3", "Page 2", "Para 12"
    location_index: int            # slide idx (pptx), page idx (pdf), para idx (docx)
    has_image_on_page: bool = False  # true if the slide/page contains an image
    has_table_on_page: bool = False  # true if the slide/page contains a table
    # Handles back to the source object so markers.py can annotate it later.
    source_handle: Optional[Any] = None


@dataclass
class Document:
    path: str
    ext: str                         # ".pptx" | ".docx" | ".pdf"
    units: List[TextUnit] = field(default_factory=list)
    # Per-location aggregated text (used for copyright long-text + citation check)
    location_text: dict = field(default_factory=dict)
    # Per-location flags: has_image, has_table
    location_flags: dict = field(default_factory=dict)
    # Raw handle to the underlying library object (python-pptx Presentation etc.)
    raw: Any = None


# ------------------------------------------------------------
# PPTX
# ------------------------------------------------------------

def iter_shapes_recursive(shapes):
    """Yield every shape under `shapes`, descending into GroupShape recursively.

    python-pptx の `slide.shapes` / GroupShape.shapes はトップレベルしか返さない
    ため、グループ化された装飾内のテキスト・表・チャート・画像が全チェッカー
    から不可視になる。この関数を経由することで、任意ネストのグループを展開
    して個々の leaf shape を訪問できる。
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        try:
            stype = shape.shape_type
        except Exception:
            stype = None
        if stype == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from iter_shapes_recursive(shape.shapes)
            except Exception:
                pass
        else:
            yield shape


def extract_pptx(path: str) -> Document:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    doc = Document(path=path, ext=".pptx", raw=prs)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        has_image = False
        has_table = False
        texts_on_slide = []

        for shape in iter_shapes_recursive(slide.shapes):
            try:
                stype = shape.shape_type
            except Exception:
                stype = None

            if stype == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
            if shape.has_table:
                has_table = True
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text or ""
                        if t.strip():
                            texts_on_slide.append(t)
                            doc.units.append(TextUnit(
                                kind="table-cell",
                                text=t,
                                location_label=f"Slide {slide_idx}",
                                location_index=slide_idx,
                                source_handle=cell,
                            ))
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs) or ""
                    if t.strip():
                        texts_on_slide.append(t)
                        doc.units.append(TextUnit(
                            kind="text",
                            text=t,
                            location_label=f"Slide {slide_idx}",
                            location_index=slide_idx,
                            source_handle=shape,
                        ))

        doc.location_flags[slide_idx] = {
            "has_image": has_image,
            "has_table": has_table,
            "label": f"Slide {slide_idx}",
        }
        doc.location_text[slide_idx] = "\n".join(texts_on_slide)

    # Propagate has_image/has_table flags to each unit
    for u in doc.units:
        flags = doc.location_flags.get(u.location_index, {})
        u.has_image_on_page = flags.get("has_image", False)
        u.has_table_on_page = flags.get("has_table", False)

    return doc


# ------------------------------------------------------------
# DOCX
# ------------------------------------------------------------

def extract_docx(path: str) -> Document:
    from docx import Document as DocxDocument

    d = DocxDocument(path)
    doc = Document(path=path, ext=".docx", raw=d)

    # Detect images / tables at document level
    has_image = False
    try:
        # python-docx exposes inline shapes
        if len(d.inline_shapes) > 0:
            has_image = True
    except Exception:
        pass
    has_table = len(d.tables) > 0

    # DOCX doesn't have natural "pages"; we treat the whole doc as location index 1
    # but track paragraph index for reporting.
    loc_idx = 1
    texts = []

    for p_idx, para in enumerate(d.paragraphs, start=1):
        t = para.text or ""
        if t.strip():
            texts.append(t)
            doc.units.append(TextUnit(
                kind="text",
                text=t,
                location_label=f"Para {p_idx}",
                location_index=loc_idx,
                source_handle=para,
            ))

    for t_idx, tbl in enumerate(d.tables, start=1):
        for r_idx, row in enumerate(tbl.rows, start=1):
            for c_idx, cell in enumerate(row.cells, start=1):
                t = cell.text or ""
                if t.strip():
                    texts.append(t)
                    doc.units.append(TextUnit(
                        kind="table-cell",
                        text=t,
                        location_label=f"Table {t_idx} R{r_idx}C{c_idx}",
                        location_index=loc_idx,
                        source_handle=cell,
                    ))

    doc.location_flags[loc_idx] = {
        "has_image": has_image,
        "has_table": has_table,
        "label": "Document",
    }
    doc.location_text[loc_idx] = "\n".join(texts)

    for u in doc.units:
        u.has_image_on_page = has_image
        u.has_table_on_page = has_table

    return doc


# ------------------------------------------------------------
# PDF
# ------------------------------------------------------------

def extract_pdf(path: str) -> Document:
    import pdfplumber

    doc = Document(path=path, ext=".pdf")

    with pdfplumber.open(path) as pdf:
        for page_idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            has_image = len(page.images) > 0
            has_table = False
            try:
                tables = page.extract_tables()
                has_table = bool(tables)
            except Exception:
                pass

            doc.location_flags[page_idx] = {
                "has_image": has_image,
                "has_table": has_table,
                "label": f"Page {page_idx}",
            }
            doc.location_text[page_idx] = text

            if text.strip():
                # Split PDF page text by paragraph (blank line) for finer reporting
                paragraphs = [p for p in text.split("\n\n") if p.strip()]
                if not paragraphs:
                    paragraphs = [text]
                for para in paragraphs:
                    doc.units.append(TextUnit(
                        kind="text",
                        text=para,
                        location_label=f"Page {page_idx}",
                        location_index=page_idx,
                        has_image_on_page=has_image,
                        has_table_on_page=has_table,
                    ))

    return doc


def extract(path: str) -> Document:
    lower = path.lower()
    if lower.endswith(".pptx"):
        return extract_pptx(path)
    if lower.endswith(".docx"):
        return extract_docx(path)
    if lower.endswith(".pdf"):
        return extract_pdf(path)
    raise ValueError(f"Unsupported file type: {path}")
